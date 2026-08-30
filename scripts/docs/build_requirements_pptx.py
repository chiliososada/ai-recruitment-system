#!/usr/bin/env python3
"""要件定義書(PowerPoint)生成スクリプト。
ケイスタンプ株式会社の AI人材採用マッチングシステム 要件定義書(である調)を
docs/requirements/AI人材採用マッチングシステム_要件定義書.pptx に生成する。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "docs/requirements/AI人材採用マッチングシステム_要件定義書.pptx"
FONT = "Yu Gothic"

NAVY = RGBColor(0x1B, 0x24, 0x55)
NAVY2 = RGBColor(0x27, 0x33, 0x6B)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
TEAL = RGBColor(0x0D, 0x94, 0x88)
ICE = RGBColor(0xEE, 0xF2, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x21, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LINE = RGBColor(0xD8, 0xDE, 0xEC)
AMBER = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW = 13.333


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, fill=None, line=None, rounded=False, line_w=1.0):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h),
    )
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def oval(s, l, t, w, h, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(s, l, t, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sa=4, sb=0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa)
        p.space_before = Pt(sb)
        for run in para:
            txt, size, bold, color = run[0], run[1], run[2], run[3]
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = FONT
            if len(run) > 4 and run[4]:
                r.font.line_spacing = run[4]
    return tb


def header(s, num, label, title):
    """Light content-slide header: left accent band, number chip, label, title."""
    rect(s, 0, 0, SW, 7.5, fill=ICE)
    rect(s, 0, 0, 0.16, 7.5, fill=INDIGO)
    chip = rect(s, 0.6, 0.5, 0.7, 0.7, fill=INDIGO, rounded=True)
    ctf = chip.text_frame
    ctf.word_wrap = False
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = num
    cr.font.size = Pt(20)
    cr.font.bold = True
    cr.font.color.rgb = WHITE
    cr.font.name = FONT
    text(s, 1.45, 0.5, 11.2, 0.4, [[(label, 12, True, TEAL)]])
    text(s, 1.45, 0.82, 11.2, 0.7, [[(title, 28, True, NAVY)]])


def footer(s, n):
    text(s, 0.6, 7.05, 9.5, 0.3,
         [[("ケイスタンプ株式会社 ｜ AI人材採用マッチングシステム 要件定義書", 9, False, MUTED)]])
    text(s, 12.0, 7.05, 0.9, 0.3, [[(str(n), 9, False, MUTED)]], align=PP_ALIGN.RIGHT)


def bullets(s, l, t, w, h, items, size=15, gap=8, color=TEXT):
    paras = []
    for it in items:
        paras.append([("•  ", size, True, INDIGO), (it, size, False, color)])
    text(s, l, t, w, h, paras, sa=gap)


# ───────────────────────── 1. 表紙 ─────────────────────────
s = slide()
rect(s, 0, 0, SW, 7.5, fill=NAVY)
rect(s, 0, 0, SW, 0.28, fill=INDIGO)
rect(s, 0, 7.22, SW, 0.28, fill=TEAL)
text(s, 1.0, 2.2, 11.3, 0.5, [[("要件定義書", 20, True, RGBColor(0xCA, 0xDC, 0xFC))]])
text(s, 1.0, 2.75, 11.3, 1.6,
     [[("AI人材採用マッチングシステム", 44, True, WHITE)]])
text(s, 1.0, 4.25, 11.3, 0.5,
     [[("AIによる人材と求人の高精度マッチング基盤", 16, False, RGBColor(0xCA, 0xDC, 0xFC))]])
rect(s, 1.0, 5.0, 4.6, 0.02, fill=RGBColor(0x44, 0x52, 0x8F))
text(s, 1.0, 5.3, 11.3, 1.2,
     [[("ケイスタンプ株式会社", 18, True, WHITE)],
      [("版数 1.0 ／ 2026-06-20", 13, False, RGBColor(0xCA, 0xDC, 0xFC))]], sa=6)

# ───────────────────────── 2. 目次 ─────────────────────────
s = slide()
rect(s, 0, 0, SW, 7.5, fill=ICE)
rect(s, 0, 0, 0.16, 7.5, fill=INDIGO)
text(s, 0.6, 0.55, 11.0, 0.7, [[("目次", 28, True, NAVY)]])
agenda = [
    ("01", "背景・課題"), ("02", "目的・ゴール"), ("03", "対象ユーザーとロール"),
    ("04", "スコープ"), ("05", "業務フロー"), ("06", "機能要件(FR-01〜FR-10)"),
    ("07", "AIマッチング要件"), ("08", "非機能要件"), ("09", "システム全体像"),
    ("10", "データ・外部連携・前提制約"),
]
for i, (num, t_) in enumerate(agenda):
    col = i // 5
    row = i % 5
    x = 0.7 + col * 6.2
    y = 1.7 + row * 1.02
    ch = oval(s, x, y, 0.6, 0.6, INDIGO if col == 0 else TEAL)
    ch.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ch.text_frame.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    rr = cp.add_run()
    rr.text = num
    rr.font.size = Pt(15)
    rr.font.bold = True
    rr.font.color.rgb = WHITE
    rr.font.name = FONT
    text(s, x + 0.8, y + 0.07, 5.2, 0.6, [[(t_, 16, True, TEXT)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 2)

# ───────────────────────── 3. 改訂履歴 ─────────────────────────
s = slide()
header(s, "—", "DOCUMENT HISTORY", "改訂履歴")
cols = [("版数", 1.4), ("日付", 2.0), ("変更区分", 2.0), ("変更内容", 5.0), ("作成者", 2.0)]
x = 0.7
rect(s, 0.7, 1.9, 12.4, 0.5, fill=NAVY)
cx = x
for name, w in cols:
    text(s, cx + 0.1, 1.97, w, 0.4, [[(name, 12, True, WHITE)]])
    cx += w
row = ["1.0", "2026-06-20", "新規作成",
       "初版。背景・目的、対象ユーザー、スコープ、業務フロー、機能要件(FR-01〜10)、非機能要件、全体像を定義した。",
       "ケイスタンプ株式会社"]
rect(s, 0.7, 2.4, 12.4, 1.05, fill=WHITE, line=LINE)
cx = x
for (name, w), val in zip(cols, row):
    text(s, cx + 0.1, 2.5, w - 0.15, 0.95, [[(val, 11, False, TEXT)]])
    cx += w
text(s, 0.7, 3.85, 12.4, 0.5,
     [[("※ 本書はケイスタンプ株式会社が、対象システムの要件を定義したものである。", 11, False, MUTED)]])
footer(s, 3)

# ───────────────────────── 4. 背景・課題 ─────────────────────────
s = slide()
header(s, "01", "BACKGROUND", "背景・課題")
issues = [
    ("ミスマッチ", "求人票と職務経歴の適合判断が定性的で、ミスマッチが生じやすい。"),
    ("選考工数", "母集団形成・スクリーニング・選考の工数が大きく、属人化している。"),
    ("説明性の不足", "適合度を客観的・説明可能に示す指標が不足している。"),
    ("多言語・保護要請", "日英中の多言語対応と、個人情報保護・権限分離の要請が高い。"),
]
for i, (h, d) in enumerate(issues):
    y = 1.85 + i * 1.18
    rect(s, 0.7, y, 12.4, 1.02, fill=WHITE, line=LINE, rounded=True)
    oval(s, 0.95, y + 0.28, 0.46, 0.46, TEAL)
    text(s, 1.75, y + 0.16, 11.1, 0.4, [[(h, 15, True, NAVY)]])
    text(s, 1.75, y + 0.52, 11.1, 0.4, [[(d, 13, False, TEXT)]])
footer(s, 4)

# ───────────────────────── 5. 目的・ゴール ─────────────────────────
s = slide()
header(s, "02", "GOALS", "目的・ゴール")
goals = [
    "履歴書からスキルを構造化抽出し、求人と候補者を双方向に高精度マッチングすることを目的とする。",
    "一致スキル・不足スキル・理由を伴う説明可能なスコアで、選考判断を支援する。",
    "行レベル権限分離による安全性と、4言語対応による利用体験の向上を実現する。",
]
bullets(s, 0.8, 1.85, 11.8, 2.2, goals, size=15, gap=12)
stats = [("10", "機能要件 (FR-01〜10)"), ("4", "対応言語 (ja/en/zh-CN/zh-TW)"), ("2", "利用ロール (求職者/企業)")]
for i, (num, lab) in enumerate(stats):
    x = 0.8 + i * 4.0
    rect(s, x, 4.5, 3.7, 1.9, fill=WHITE, line=LINE, rounded=True)
    text(s, x, 4.7, 3.7, 1.0, [[(num, 54, True, INDIGO)]], align=PP_ALIGN.CENTER)
    text(s, x + 0.15, 5.85, 3.4, 0.5, [[(lab, 12, True, TEXT)]], align=PP_ALIGN.CENTER)
footer(s, 5)

# ───────────────────────── 6. 対象ユーザーとロール ─────────────────────────
s = slide()
header(s, "03", "USERS & ROLES", "対象ユーザーとロール")
roles = [
    ("求職者", "job_seeker", "履歴書登録、AIスキル分析、推薦閲覧、応募、メッセージ、面接対応を行う。", INDIGO),
    ("企業担当者", "company_member", "企業・求人管理、人材検索、選考、面接提案、採用判断を行う。", TEAL),
    ("未認証", "anonymous", "公開求人・公開企業の閲覧のみを行う。", NAVY2),
    ("運用", "service", "解析等の内部処理を担う信頼経路である(権限を限定)。", MUTED),
]
for i, (h, code, d, c) in enumerate(roles):
    col = i % 2
    rw = i // 2
    x = 0.7 + col * 6.25
    y = 1.9 + rw * 2.35
    rect(s, x, y, 5.95, 2.1, fill=WHITE, line=LINE, rounded=True)
    rect(s, x, y, 0.14, 2.1, fill=c, rounded=False)
    text(s, x + 0.35, y + 0.25, 5.4, 0.5, [[(h, 18, True, NAVY), ("　" + code, 11, False, MUTED)]])
    text(s, x + 0.35, y + 0.9, 5.4, 1.0, [[(d, 13, False, TEXT)]])
footer(s, 6)

# ───────────────────────── 7. スコープ ─────────────────────────
s = slide()
header(s, "04", "SCOPE", "スコープ")
rect(s, 0.7, 1.9, 6.0, 4.7, fill=WHITE, line=LINE, rounded=True)
rect(s, 0.7, 1.9, 6.0, 0.6, fill=INDIGO, rounded=True)
text(s, 0.9, 2.0, 5.6, 0.4, [[("対象(In Scope)", 15, True, WHITE)]])
bullets(s, 0.95, 2.7, 5.5, 3.7, [
    "FR-01 認証・ユーザー管理", "FR-02 履歴書アップロード・解析", "FR-03 AIスキル分析",
    "FR-04 企業・求人管理", "FR-05 AIマッチング", "FR-06 人材検索・候補者詳細",
    "FR-07 企業・求人検索", "FR-08 メッセージ・通知", "FR-09 多言語対応",
    "FR-10 候補者比較・面接・採用フロー",
], size=12.5, gap=5)
rect(s, 7.0, 1.9, 5.9, 4.7, fill=WHITE, line=LINE, rounded=True)
rect(s, 7.0, 1.9, 5.9, 0.6, fill=NAVY2, rounded=True)
text(s, 7.2, 2.0, 5.5, 0.4, [[("対象外(Out of Scope)", 15, True, WHITE)]])
bullets(s, 7.25, 2.7, 5.4, 3.7, [
    "課金・決済、給与計算", "外部求人媒体との連携(本書時点)", "SNSログイン連携",
    "ネイティブモバイルアプリ", "本番環境への実デプロイ(別途承認のうえ実施)",
], size=13, gap=10, color=TEXT)
text(s, 7.25, 5.9, 5.4, 0.6, [[("※ 対象外項目は将来拡張として別途検討する。", 11, False, MUTED)]])
footer(s, 7)

# ───────────────────────── 8. 業務フロー ─────────────────────────
s = slide()
header(s, "05", "BUSINESS FLOW", "業務フロー")


def flow(s, y, title, color, steps):
    text(s, 0.7, y, 12.0, 0.4, [[(title, 15, True, color)]])
    n = len(steps)
    bw = 1.95
    gap = (12.4 - bw * n) / (n - 1) if n > 1 else 0
    for i, st in enumerate(steps):
        x = 0.7 + i * (bw + gap)
        rect(s, x, y + 0.5, bw, 0.95, fill=WHITE, line=color, rounded=True, line_w=1.5)
        text(s, x + 0.05, y + 0.5, bw - 0.1, 0.95,
             [[(st, 12.5, True, NAVY)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            text(s, x + bw - 0.02, y + 0.5, gap + 0.04, 0.95,
                 [[("→", 18, True, color)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


flow(s, 2.0, "求職者フロー", INDIGO,
     ["登録・確認", "履歴書登録", "AIスキル分析", "推薦・応募", "面接・採用"])
flow(s, 4.4, "企業フロー", TEAL,
     ["企業・求人作成", "人材検索", "候補者選考", "面接提案", "採用判断"])
footer(s, 8)

# ───────────────────────── 9-10. 機能要件 ─────────────────────────
FR = [
    ("01", "認証・ユーザー管理", "登録・ログイン・メール確認・アカウント更新。求職者/企業の二ロールを管理する。"),
    ("02", "履歴書アップロード・解析", "PDF/DOCXを検証・ウイルススキャンし、テキスト抽出と解析を非同期に行う。"),
    ("03", "AIスキル分析", "スキル・熟練度・経験年数・キャリア提案・推奨学習を抽出し、可視化する。"),
    ("04", "企業・求人管理", "企業の作成・更新、求人のCRUDと公開範囲(下書き/公開・公開/非公開)を管理する。"),
    ("05", "AIマッチング", "ベクトル近傍探索とルール加重で、説明可能な適合度スコアを算定する。"),
    ("06", "人材検索・候補者詳細", "企業が候補者を検索・絞り込み・比較し、関係成立時のみ連絡先等を開示する。"),
    ("07", "企業・求人検索", "公開求人・公開企業の検索・絞り込み・ページングを未認証でも提供する。"),
    ("08", "メッセージ・通知", "スレッド型メッセージ、未読管理、応募・面接等のイベント通知を提供する。"),
    ("09", "多言語対応", "ja/en/zh-CN/zh-TW の4言語でUIとエラーを提供し、辞書整合を担保する。"),
    ("10", "候補者比較・面接・採用フロー", "応募の段階遷移、ショートリスト、候補者比較、面接提案・応答を管理する。"),
]


def fr_slide(s, title_sub, items, pageno):
    header(s, "06", "FUNCTIONAL REQUIREMENTS", "機能要件 " + title_sub)
    for i, (num, h, d) in enumerate(items):
        y = 1.85 + i * 0.95
        rect(s, 0.7, y, 12.4, 0.82, fill=WHITE, line=LINE, rounded=True)
        chip = rect(s, 0.92, y + 0.16, 0.95, 0.5, fill=INDIGO, rounded=True)
        chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = chip.text_frame.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        rr = cp.add_run()
        rr.text = "FR-" + num
        rr.font.size = Pt(12)
        rr.font.bold = True
        rr.font.color.rgb = WHITE
        rr.font.name = FONT
        text(s, 2.1, y + 0.1, 10.8, 0.35, [[(h, 14.5, True, NAVY)]])
        text(s, 2.1, y + 0.45, 10.8, 0.32, [[(d, 11.5, False, TEXT)]])
    footer(s, pageno)


fr_slide(slide(), "(FR-01〜FR-05)", FR[:5], 9)
fr_slide(slide(), "(FR-06〜FR-10)", FR[5:], 10)

# ───────────────────────── 11. AIマッチング要件 ─────────────────────────
s = slide()
header(s, "07", "AI MATCHING", "AIマッチング要件")
bullets(s, 0.8, 1.85, 11.9, 3.2, [
    "履歴書テキストから、スキル・経験年数・キャリア提案を構造化抽出する(スキーマ検証を行う)。",
    "候補者・求人を384次元のベクトルへ符号化し、pgvectorの近傍探索で母集団を抽出する。",
    "スキル・経験・給与・勤務地・言語・ベクトル類似度の加重により、0〜100の適合度を算定する。",
    "スコアは決定論的かつ版管理され、一致/不足スキルと理由を提示する(説明可能性を担保する)。",
    "外部LLM障害時も、回路遮断・タイムアウト・再試行により処理の安定性を確保する。",
], size=14.5, gap=12)
band = [("入力", "履歴書・プロフィール・求人"), ("解析", "スキル構造化(検証付き)"),
        ("符号化", "384次元ベクトル化"), ("算定", "近傍探索＋加重スコア"), ("出力", "適合度・理由・一致/不足")]
n = len(band)
bw = 2.3
gap = (12.4 - bw * n) / (n - 1)
for i, (h, d) in enumerate(band):
    x = 0.7 + i * (bw + gap)
    rect(s, x, 5.4, bw, 1.2, fill=NAVY if i % 2 == 0 else NAVY2, rounded=True)
    text(s, x + 0.1, 5.55, bw - 0.2, 0.4, [[(h, 13, True, WHITE)]], align=PP_ALIGN.CENTER)
    text(s, x + 0.1, 5.95, bw - 0.2, 0.6, [[(d, 10, False, RGBColor(0xCA, 0xDC, 0xFC))]], align=PP_ALIGN.CENTER)
    if i < n - 1:
        text(s, x + bw - 0.05, 5.4, gap + 0.1, 1.2, [[("›", 16, True, TEAL)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 11)

# ───────────────────────── 12. 非機能要件 ─────────────────────────
s = slide()
header(s, "08", "NON-FUNCTIONAL", "非機能要件")
nfr = [
    ("セキュリティ", "行レベル権限分離(RLS)、厳格なHTTPヘッダ、経路別レート制限を行う。"),
    ("プライバシー", "個人情報を最小化し、ログ秘匿と所有者限定アクセスを徹底する。"),
    ("性能", "主要APIの低レイテンシと、フロント初期JSの容量予算を満たす。"),
    ("可用性・信頼性", "障害耐性(回路遮断/再試行)、耐久ジョブ処理、安全な停止を行う。"),
    ("多言語", "ja/en/zh-CN/zh-TW を全UIで提供する。"),
    ("アクセシビリティ", "WCAG 2.1 AA に準拠する。"),
    ("可観測性・運用", "死活/受入監視、構造化ログ、相関IDを備える。"),
    ("拡張性", "ランタイム抽象化とアダプタにより、提供基盤を切替可能とする。"),
]
for i, (h, d) in enumerate(nfr):
    col = i % 2
    rw = i // 2
    x = 0.7 + col * 6.25
    y = 1.85 + rw * 1.22
    rect(s, x, y, 5.95, 1.08, fill=WHITE, line=LINE, rounded=True)
    rect(s, x, y, 0.12, 1.08, fill=TEAL)
    text(s, x + 0.3, y + 0.13, 5.5, 0.35, [[(h, 14, True, NAVY)]])
    text(s, x + 0.3, y + 0.5, 5.5, 0.5, [[(d, 11.5, False, TEXT)]])
footer(s, 12)

# ───────────────────────── 13. システム全体像 ─────────────────────────
s = slide()
header(s, "09", "ARCHITECTURE", "システム全体像")
layers = [
    ("フロントエンド", "React SPA。多言語・レスポンシブ・デザインシステム。", INDIGO),
    ("API", "Fastify。入出力検証、リクエスト毎の権限コンテキスト、API契約公開。", NAVY2),
    ("データ", "PostgreSQL + pgvector。全テーブルで行レベル権限分離(RLS)。", NAVY),
    ("AI / 連携", "LLM・埋め込みアダプタ(外部プロバイダ切替)。決定論モックで検証可能。", TEAL),
]
for i, (h, d, c) in enumerate(layers):
    y = 1.9 + i * 1.0
    rect(s, 0.7, y, 8.4, 0.86, fill=c, rounded=True)
    text(s, 0.95, y + 0.12, 2.4, 0.6, [[(h, 15, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.3, y + 0.12, 5.6, 0.62, [[(d, 11.5, False, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
rect(s, 9.4, 1.9, 3.5, 3.96, fill=WHITE, line=LINE, rounded=True)
text(s, 9.6, 2.05, 3.1, 0.4, [[("2系統ランタイム", 13, True, NAVY)]])
bullets(s, 9.62, 2.55, 3.1, 3.2, [
    "検証用: 内蔵DB+決定論モック(資格情報不要)",
    "本番: マネージドDB・認証・ストレージ・外部LLM",
    "同一コードで両系統を切替",
], size=11, gap=8)
text(s, 0.7, 6.1, 12.4, 0.7,
     [[("認証基盤・オブジェクトストレージ・ウイルススキャンを外部サービスとして連携する。", 12, False, MUTED)]])
footer(s, 13)

# ───────────────────────── 14. データ・外部連携・前提制約 ─────────────────────────
s = slide()
header(s, "10", "DATA / INTEGRATION / CONSTRAINTS", "データ・外部連携・前提制約")
rect(s, 0.7, 1.9, 6.0, 2.3, fill=WHITE, line=LINE, rounded=True)
text(s, 0.95, 2.05, 5.5, 0.4, [[("主要データ", 14, True, NAVY)]])
text(s, 0.95, 2.5, 5.5, 1.6,
     [[("ユーザー／候補者／履歴書／スキル分析／企業／求人／マッチング結果／応募／面接／会話／通知 を中核エンティティとする。",
        12, False, TEXT)]])
rect(s, 7.0, 1.9, 5.9, 2.3, fill=WHITE, line=LINE, rounded=True)
text(s, 7.25, 2.05, 5.4, 0.4, [[("外部連携", 14, True, NAVY)]])
bullets(s, 7.27, 2.5, 5.4, 1.6, [
    "外部LLM(Anthropic / OpenAI、設定で切替)",
    "埋め込み生成・ベクトル検索(pgvector)",
    "オブジェクトストレージ・認証基盤・ウイルススキャン(ClamAV)",
], size=11.5, gap=6)
rect(s, 0.7, 4.4, 12.2, 2.2, fill=WHITE, line=LINE, rounded=True)
text(s, 0.95, 4.55, 11.7, 0.4, [[("前提・制約", 14, True, NAVY)]])
bullets(s, 0.97, 5.0, 11.7, 1.5, [
    "個人情報保護法等の関連法令を遵守する。",
    "データベース移行は追加専用(前進専用)とし、API・DTOは後方互換を維持する。",
    "本番環境への実デプロイ・課金リソースの作成は、別途承認のうえ実施する。",
], size=12.5, gap=7)
footer(s, 14)

# ───────────────────────── 15. まとめ ─────────────────────────
s = slide()
rect(s, 0, 0, SW, 7.5, fill=NAVY)
rect(s, 0, 0, SW, 0.28, fill=INDIGO)
rect(s, 0, 7.22, SW, 0.28, fill=TEAL)
text(s, 1.0, 2.3, 11.3, 0.9, [[("まとめ", 36, True, WHITE)]])
bullets_paras = [
    "本要件定義は、FR-01〜FR-10 と非機能要件を定め、説明可能なAIマッチングと安全な権限分離を中核とする。",
    "4言語対応とアクセシビリティを満たし、求職者・企業双方の体験向上を目指す。",
]
paras = [[("•  ", 16, True, TEAL), (t_, 16, False, RGBColor(0xE6, 0xEC, 0xFB))] for t_ in bullets_paras]
text(s, 1.0, 3.5, 11.3, 2.0, paras, sa=14)
text(s, 1.0, 6.3, 11.3, 0.5, [[("ケイスタンプ株式会社", 15, True, WHITE)]])

prs.save(OUT)
import os
print("SLIDES:", len(prs.slides._sldIdLst))
print(f"SAVED: {OUT} ({os.path.getsize(OUT)} bytes)")
